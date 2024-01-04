# Author: Barış Özdizdar
# Transforming the pdf into a string of text and saving it into a txt file.
import os
import openai
from fastapi import UploadFile
import fitz
import pandas as pd
from docx import Document
from io import BytesIO
from pptx import Presentation
import chromadb

class Transformer:
    
    openai.api_key = os.environ["OPENAI_API_KEY"]

    PROMPT = "Please write a detailed summary of the following:"
    # MAIN_PATH = "/Users/orbiszeus/Master_AI_LlamaIndex-1"
    MAIN_PATH = os.
    
    #This method saves uploadede pdfs into a txt file inside the curr directory.
    def save_file(content, filepath):
        with open(filepath, 'w', encoding='utf-8') as outfile:
            outfile.write(content)
            
    #This method parses the uploaded binary pdf into a dict of strings later,
    #splitted into correct form whereas the original pdf looks like into a string of text.
    def parse_files(self, files : list[UploadFile], unique_id, file_extension_list):
        raw_texts = []
        for i in range(len(file_extension_list)):
            match file_extension_list[i]:
                case "pdf":
                    raw_texts.append(self.parse_pdf(files[i]))
                case "docx":
                    raw_texts.append(self.parse_docx(files[i]))
                case "pptx":
                    raw_texts.append(self.parse_pptx(files[i]))
                case "csv":
                    raw_texts.append(self.parse_csv(files[i]))
                case "xlsx":    
                    raw_texts.append(self.parse_xlsx(files[i]))
                case _:
                    return "File extension is not supported."
              
        dir_for_creating_index = Transformer.create_folder_system(unique_id=unique_id, raw_text=raw_texts)
        return [{"texts": raw_texts, "uniqueId" : str(unique_id)}, raw_texts, dir_for_creating_index]
    
    def create_folder_system(unique_id, raw_text):
        folder_name = str("TestData_" + unique_id)
        os.makedirs(str(os.path.join(Transformer.MAIN_PATH, folder_name)), exist_ok=True)
        dir_for_creating_index = "./" + folder_name
        for i, text in enumerate(raw_text):
            author = 'data_' + unique_id + f'_{i}.txt'
            dynamic_path_name = str(os.path.join(Transformer.MAIN_PATH, folder_name, author))
            Transformer.save_file(text, dynamic_path_name)    
        return dir_for_creating_index
    
    def parse_pdf(self, files):
        raw_texts = ''
        file_bytes = files.file.read()
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page in doc:
            raw_texts += (page.get_text())
        return raw_texts
    
    def parse_docx(self, files):
        raw_texts = ''
        # Read the file into bytes
        file_bytes = files.file.read()
        # Open the file with python-docx
        doc = Document(BytesIO(file_bytes))

        # Extract the text
        for paragraph in doc.paragraphs:
            raw_texts += paragraph.text + '\n'

        return raw_texts
    
    def parse_pptx(self, files):
        raw_texts = ''

        # Read the file into bytes
        file_bytes = files.file.read()

        # Open the file with python-pptx
        prs = Presentation(BytesIO(file_bytes))

        # Extract the text
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            raw_texts += run.text + '\n'

        return raw_texts
    
    
    def parse_csv(self, files):
        raw_texts = ''
        # Read the file into bytes
        file_bytes = files.file.read()
        # Open the file with pandas
        df = pd.read_csv(BytesIO(file_bytes))
        # Convert the DataFrame to a CSV string and strip the trailing newline
        raw_texts = df.to_csv(index=False, lineterminator='\n').strip()

        return raw_texts
    
    
    def parse_xlsx(self, files):
        raw_texts = ''
    # Read the file into bytes
        file_bytes = files.file.read()
        # Open the file with pandas
        df = pd.read_excel(BytesIO(file_bytes))
        # Convert the DataFrame to a CSV string and strip the trailing newline
        raw_texts = df.to_csv(index=False, lineterminator='\n').strip()

        return raw_texts